import os
from typing import List, Dict, Any
import scrapy
from scrapy.http import Request, FormRequest
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import pylibmagic
import magic
import requests
from pydantic import BaseModel, Field
from dotenv import load_dotenv

load_dotenv()

class KMSItem(BaseModel):
    title: str = Field(description="文档标题")
    content: str = Field(description="文档内容")
    attachments: List[Dict[str, Any]] = Field(default_factory=list, description="附件信息")

class ConfluenceSpider(scrapy.Spider):
    name = 'confluence'
    custom_settings = {
        'DOWNLOAD_DELAY': 1,
        'COOKIES_ENABLED': True,
        'CONCURRENT_REQUESTS': 1
    }

    def __init__(self, *args, **kwargs):
        super().__init__(*args, **kwargs)
        self.start_urls = [kwargs.get('start_url', 'http://kms.new-see.com:8090')]
        self.basic_auth = ('newsee', 'newsee')
        self.auth = {
            'os_username': 'zengdi',
            'os_password': '808611'
        }

        self.baichuan_config = {
            'api_key': os.getenv('BAI_CH_API_KEK'),
            'api_url': 'https://api.baichuan-ai.com/v1/chat/completions'
        }

    def start_requests(self):
        for url in self.start_urls:
            yield Request(
                url,
                callback=self.parse,
                meta={
                    'cookiejar': 1,
                    'dont_redirect': False
                },
                dont_filter=True
            )

    def parse(self, response):
        # 处理登录
        if response.css('#loginform'):
            return self.login(response)

        # 解析导航树
        tree_links = response.css('.plugin-tabmeta-details a::attr(href)').getall()
        for link in tree_links:
            if 'pageId' in link:
                yield response.follow(link, callback=self.parse_content)

    def login(self, response):
        print("🚀 ~ self, response:", self, response)

        res = FormRequest.from_response(
            response,
            formcss='#loginform',
            formdata={'os_username': self.auth['os_username'], 'os_password': self.auth['os_password']},
            clickdata={'css': '#loginButton'},
            callback=self.after_login,
            dont_filter=True,
            meta={'cookies': response.meta.get('cookies')}
        )
        print("🚀 ~ self, response:", self, res)
        return res

    def after_login(self, response):
        if 'login' not in response.url:
            return self.parse(response)
        else:
            self.logger.error('登录失败')
            return None

    def optimize_content(self, content: str) -> str:
        headers = {
            'Content-Type': 'application/json',
            'Authorization': f'Bearer {self.baichuan_config["api_key"]}'
        }

        data = {
            'model': 'Baichuan4',
            'messages': [
                {
                    'role': 'system',
                    'content': '你是一个专业的文档优化助手，需要对输入的文档内容进行结构化和优化处理，使其更加清晰易读，同时保持原有的核心信息和专业性。'
                },
                {
                    'role': 'user',
                    'content': content
                }
            ],
            'temperature': 0.3,
            'stream': False
        }

        try:
            response = requests.post(self.baichuan_config['api_url'], headers=headers, json=data)
            response.raise_for_status()
            result = response.json()
            return result['choices'][0]['message']['content']
        except Exception as e:
            self.logger.error(f'百川API调用失败: {str(e)}')
            return content

    def parse_content(self, response):
        soup = BeautifulSoup(response.text, 'html.parser')
        title = soup.select_one('#title-text').get_text(strip=True)
        content = soup.select_one('#main-content')

        # 处理附件
        attachments = []
        for attachment in soup.select('.attachment-content'):
            file_url = attachment.select_one('a::attr(href)').get()
            file_type = magic.from_file(file_url, mime=True)

            if 'image' in file_type:
                text = self.process_image(file_url)
            elif 'pdf' in file_type:
                text = self.process_pdf(file_url)
            elif 'word' in file_type:
                text = self.process_word(file_url)
            elif 'powerpoint' in file_type:
                text = self.process_ppt(file_url)
            else:
                text = None

            if text:
                attachments.append({
                    'url': file_url,
                    'type': file_type,
                    'content': text
                })

        # 使用百川API优化内容
        optimized_content = self.optimize_content(content.get_text())

        yield KMSItem(
            title=title,
            content=optimized_content,
            attachments=attachments
        )

    @staticmethod
    def process_image(image_path):
        image = Image.open(image_path)
        return pytesseract.image_to_string(image, lang='chi_sim')

    @staticmethod
    def process_pdf(pdf_path):
        pages = convert_from_path(pdf_path)
        text = ''
        for page in pages:
            text += pytesseract.image_to_string(page, lang='chi_sim')
        return text

    @staticmethod
    def process_word(docx_path):
        doc = Document(docx_path)
        return '\n'.join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    def process_ppt(pptx_path):
        prs = Presentation(pptx_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + '\n'
        return text