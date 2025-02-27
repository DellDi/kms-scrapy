import os
import mimetypes
import urllib.parse
import logging
from typing import Dict, Any, List, Optional

from bs4 import BeautifulSoup, Tag
from PIL import Image
from docx import Document
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
import magic

from pydantic import BaseModel, Field


class KMSItem(BaseModel):
    """KMS文档项目"""

    title: str
    depth_info: Optional[dict] = Field(default=None, description="页面深度信息")
    content: str
    attachments: Optional[List[dict]] = []

    model_config = {
        "arbitrary_types_allowed": True,
        "json_schema_extra": {
            "examples": [{"title": "示例文档", "content": "文档内容", "attachments": []}]
        },
    }


class ContentParser:
    """内容解析器，处理页面内容和附件"""

    def __init__(
        self, enable_text_extraction: bool = True, content_optimizer=None, auth_manager=None
    ):
        """初始化解析器

        Args:
            enable_text_extraction: 是否启用对应的附件的文本提取功能，默认为False
            content_optimizer: 内容优化器实例，用于美化提取的文本
            auth_manager: 认证管理器实例，用于处理认证请求
        """
        self.enable_text_extraction = enable_text_extraction
        self.content_optimizer = content_optimizer
        self.auth_manager = auth_manager
        self.logger = logging.getLogger(__name__)
        self._spider_callback = None  # 存储Spider提供的回调函数

    def set_callback(self, callback):
        """设置Spider提供的回调函数

        Args:
            callback: Spider中处理附件下载的回调函数
        """
        self._spider_callback = callback
        self.logger.info("已设置Spider回调函数")

    @staticmethod
    def parse_page_content(html_content: str) -> tuple[Tag, Tag]:
        """解析页面内容，返回标题和正文"""
        soup = BeautifulSoup(html_content, "html.parser")
        titleDom = soup.select_one("#title-text")
        contentDom = soup.select_one("#main-content")
        # linkUrl = soup.select_one("#title-text a")["href"]
        return titleDom, contentDom

    @staticmethod
    def process_image(image_path: str) -> str:
        """处理图片文件，提取文本"""
        try:
            image = Image.open(image_path)
            return pytesseract.image_to_string(image, lang="chi_sim")
        except pytesseract.TesseractNotFoundError:
            logging.error(
                "Tesseract OCR未安装或未添加到PATH中。请参考README文件安装必要的系统依赖。"
            )
            return None
        except Exception as e:
            logging.warning(f"图片文本提取失败: {str(e)}")
            return None

    @staticmethod
    def process_pdf(pdf_path: str) -> str:
        """处理PDF文件，提取文本"""
        try:
            pages = convert_from_path(pdf_path)
            text = ""
            for page in pages:
                text += pytesseract.image_to_string(page, lang="chi_sim")
            return text
        except (pytesseract.TesseractNotFoundError, Exception) as e:
            logging.warning(f"PDF文本提取失败: {str(e)}")
            return None

    @staticmethod
    def process_word(docx_path: str) -> str:
        """处理Word文件，提取文本"""
        try:
            doc = Document(docx_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            logging.warning(f"Word文本提取失败: {str(e)}")
            return None

    @staticmethod
    def process_ppt(pptx_path: str) -> str:
        """处理PPT文件，提取文本"""
        try:
            prs = Presentation(pptx_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            logging.warning(f"PPT文本提取失败: {str(e)}")
            return None

    def handle_downloaded_file(self, response):
        """处理下载完成的文件响应"""
        self.logger.info(f"收到文件下载响应: {response.url}, 状态码: {response.status}")

        if response.status != 200:
            self.logger.error(f"附件下载失败: {response.url}, 状态码: {response.status}")
            return None

        temp_path = None
        try:
            # 获取文件类型
            file_type = magic.from_buffer(response.body, mime=True)

            # 处理文件名
            file_name = os.path.basename(response.url).split("?")[0]
            file_name = urllib.parse.unquote(file_name)

            # 处理文件后缀
            if "." not in file_name:
                content_type = response.headers.get("Content-Type", "")
                ext = mimetypes.guess_extension(content_type)
                if ext:
                    file_name = f"{file_name}{ext}"
                    
            # 检查是否需要过滤此附件（基于实际MIME类型和文件大小）
            attachment_filters = response.meta.get("attachment_filters")
            if attachment_filters and attachment_filters.get("enabled", False):
                # 检查MIME类型
                excluded_mime_types = attachment_filters.get("excluded_mime_types", [])
                if any(file_type.startswith(excluded) for excluded in excluded_mime_types):
                    self.logger.info(f"附件 {file_name} 因实际MIME类型 {file_type} 被过滤")
                    return None
                
                # 检查文件大小
                max_size_mb = attachment_filters.get("max_size_mb", 50)
                file_size_mb = len(response.body) / (1024 * 1024)  # 转换为MB
                if file_size_mb > max_size_mb:
                    self.logger.info(f"附件 {file_name} 因大小 {file_size_mb:.2f}MB 超过限制 {max_size_mb}MB 被过滤")
                    return None

            # 创建临时文件
            temp_path = f"temp_{file_name}"
            with open(temp_path, "wb") as f:
                f.write(response.body)

            # 处理文本提取
            text = None
            if self.enable_text_extraction:
                try:
                    if "image" in file_type:
                        text = self.process_image(temp_path)
                    elif "pdf" in file_type:
                        text = self.process_pdf(temp_path)
                    elif "word" in file_type:
                        text = self.process_word(temp_path)
                    elif "powerpoint" in file_type:
                        text = self.process_ppt(temp_path)

                    if text and self.content_optimizer:
                        text = self.content_optimizer.optimize(content=text, spiderUrl=response.url)
                        file_type = "text/markdown"
                except Exception as e:
                    self.logger.error(f"文本提取失败: {str(e)}")

            result = {
                "url": response.url,
                "type": file_type,
                "filename": file_name,
                "content": response.body,
                "extracted_text": text,
            }

            self.logger.info(f"附件处理完成: {file_name}")
            return result

        except Exception as e:
            self.logger.error(f"附件处理失败: {str(e)}")
            return None

        finally:
            if temp_path and os.path.exists(temp_path):
                try:
                    os.remove(temp_path)
                except OSError as e:
                    self.logger.error(f"清理临时文件失败: {str(e)}")

    def process_attachment(self, file_url: str):
        """处理附件

        Args:
            file_url: 附件URL

        Returns:
            Request: Scrapy请求对象，将由Spider执行
        """
        if not self.auth_manager:
            self.logger.error("未提供auth_manager，无法处理附件下载")
            return None

        self.logger.info(f"开始处理附件下载: {file_url}")

        # 检查是否需要过滤此附件
        from crawler.core.config import config
        
        if config.spider.attachment_filters.get("enabled", False):
            # 1. 检查文件扩展名
            file_name = os.path.basename(file_url).split("?")[0]
            file_name = urllib.parse.unquote(file_name)
            file_ext = os.path.splitext(file_name)[1].lower()
            
            # 检查扩展名过滤
            excluded_extensions = config.spider.attachment_filters.get("excluded_extensions", [])
            if file_ext in excluded_extensions:
                self.logger.info(f"附件 {file_name} 因扩展名 {file_ext} 被过滤")
                return None
                
            # 2. 检查URL中的MIME类型提示（如果有）
            mime_hint = None
            if "content-type=" in file_url.lower():
                mime_parts = [p for p in file_url.split("&") if "content-type=" in p.lower()]
                if mime_parts:
                    mime_hint = mime_parts[0].split("=")[1]
                    
            # 如果URL中有MIME类型提示，检查是否在排除列表中
            excluded_mime_types = config.spider.attachment_filters.get("excluded_mime_types", [])
            if mime_hint and any(excluded in mime_hint.lower() for excluded in excluded_mime_types):
                self.logger.info(f"附件 {file_name} 因MIME类型提示 {mime_hint} 被过滤")
                return None

        # 创建下载请求
        request = self.auth_manager.create_authenticated_request(
            url=file_url,
            # 默认使用handle_downloaded_file，但允许外部覆盖
            callback=self._spider_callback or self.handle_downloaded_file,
            meta={
                "handle_httpstatus_list": [200],
                "download_timeout": 180,  # 延长下载超时时间
                "dont_retry": False,
                "dont_merge_cookies": False,
                "is_attachment": True,  # 标记这是附件下载请求
                "attachment_filters": config.spider.attachment_filters if config.spider.attachment_filters.get("enabled", False) else None,
                # 将过滤配置传递给回调函数，避免重复导入
            },
        )

        return request  # 返回请求对象，由Spider处理执行
